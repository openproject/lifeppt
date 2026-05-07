from __future__ import annotations

from math import cos, pi, sin
from pathlib import Path
import random

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "assets"
OUT_FILE = ROOT / "小学生手卫生.pptx"

W, H = 1920, 1080
SS = 2
SLIDE_W, SLIDE_H = 13.333333, 7.5

NAVY = (27, 61, 96)
BLUE = (39, 151, 227)
CYAN = (71, 204, 226)
MINT = (111, 222, 177)
GOLD = (255, 205, 92)
CORAL = (255, 126, 112)
PINK = (255, 151, 184)
LAVENDER = (157, 142, 255)
CREAM = (255, 249, 232)
WHITE = (255, 255, 255)
SKIN = (255, 205, 165)
SKIN_DARK = (226, 154, 119)
SOAP = (127, 219, 230)
SHADOW = (39, 73, 105, 70)


def sc(v: float | int) -> int:
    return int(round(v * SS))


def sbox(box):
    return tuple(sc(v) for v in box)


def rgba(color, alpha=255):
    return (*color, alpha) if len(color) == 3 else color


def mix(a, b, t):
    return tuple(int(a[i] * (1 - t) + b[i] * t) for i in range(3))


class Art:
    def __init__(self, top, bottom, seed=0):
        self.seed = seed
        self.img = Image.new("RGBA", (W * SS, H * SS), rgba(top))
        self.draw = ImageDraw.Draw(self.img, "RGBA")
        for y in range(H * SS):
            t = y / (H * SS - 1)
            self.draw.line((0, y, W * SS, y), fill=rgba(mix(top, bottom, t)))
        self._aurora()
        self._pattern()

    def _aurora(self):
        random.seed(self.seed)
        colors = [CYAN, MINT, LAVENDER, PINK, GOLD]
        layer = Image.new("RGBA", self.img.size, (0, 0, 0, 0))
        d = ImageDraw.Draw(layer, "RGBA")
        for _ in range(8):
            x = random.randint(-250, W - 150)
            y = random.randint(-160, H - 120)
            rw = random.randint(360, 720)
            rh = random.randint(250, 520)
            c = random.choice(colors)
            d.ellipse(sbox((x, y, x + rw, y + rh)), fill=rgba(c, random.randint(32, 64)))
        layer = layer.filter(ImageFilter.GaussianBlur(sc(45)))
        self.img.alpha_composite(layer)
        self.draw = ImageDraw.Draw(self.img, "RGBA")

    def _pattern(self):
        random.seed(self.seed + 100)
        for _ in range(80):
            x = random.randint(20, W - 20)
            y = random.randint(20, H - 20)
            r = random.randint(2, 5)
            self.draw.ellipse(sbox((x - r, y - r, x + r, y + r)), fill=(255, 255, 255, random.randint(45, 95)))
        for _ in range(16):
            x = random.randint(40, W - 40)
            y = random.randint(40, H - 40)
            r = random.randint(18, 44)
            self.draw.ellipse(sbox((x - r, y - r, x + r, y + r)), outline=(255, 255, 255, 55), width=sc(2))

    def shadow_round(self, box, radius=36, blur=22, alpha=80, offset=(0, 14)):
        layer = Image.new("RGBA", self.img.size, (0, 0, 0, 0))
        d = ImageDraw.Draw(layer, "RGBA")
        moved = (box[0] + offset[0], box[1] + offset[1], box[2] + offset[0], box[3] + offset[1])
        d.rounded_rectangle(sbox(moved), radius=sc(radius), fill=(20, 61, 96, alpha))
        layer = layer.filter(ImageFilter.GaussianBlur(sc(blur)))
        self.img.alpha_composite(layer)
        self.draw = ImageDraw.Draw(self.img, "RGBA")

    def rounded(self, box, radius=36, fill=(255, 255, 255, 165), outline=(255, 255, 255, 170), width=2, shadow=True):
        if shadow:
            self.shadow_round(box, radius=radius)
        self.draw.rounded_rectangle(sbox(box), radius=sc(radius), fill=fill, outline=outline, width=sc(width))
        x1, y1, x2, y2 = box
        self.draw.line(sbox((x1 + radius, y1 + 8, x2 - radius, y1 + 8)), fill=(255, 255, 255, 110), width=sc(2))

    def circle(self, cx, cy, r, fill, outline=(255, 255, 255, 180), width=3, shadow=True):
        if shadow:
            layer = Image.new("RGBA", self.img.size, (0, 0, 0, 0))
            d = ImageDraw.Draw(layer, "RGBA")
            d.ellipse(sbox((cx - r, cy - r + 12, cx + r, cy + r + 12)), fill=SHADOW)
            layer = layer.filter(ImageFilter.GaussianBlur(sc(20)))
            self.img.alpha_composite(layer)
            self.draw = ImageDraw.Draw(self.img, "RGBA")
        self.draw.ellipse(sbox((cx - r, cy - r, cx + r, cy + r)), fill=fill, outline=outline, width=sc(width))

    def save(self, path):
        out = self.img.resize((W, H), Image.Resampling.LANCZOS).convert("RGB")
        out.save(path, quality=95)


def star_points(cx, cy, outer, inner, points=5):
    pts = []
    for i in range(points * 2):
        angle = -pi / 2 + i * pi / points
        r = outer if i % 2 == 0 else inner
        pts.append((cx + cos(angle) * r, cy + sin(angle) * r))
    return pts


def draw_sparkle(a: Art, cx, cy, r, color=GOLD):
    a.draw.polygon([tuple(sc(v) for v in p) for p in star_points(cx, cy, r, r * 0.42)], fill=rgba(color, 225), outline=rgba(WHITE, 190))


def draw_drop(a: Art, cx, cy, s, color=CYAN, alpha=245):
    pts = [
        (cx, cy - 1.25 * s),
        (cx - .72 * s, cy - .10 * s),
        (cx - .54 * s, cy + .62 * s),
        (cx, cy + .95 * s),
        (cx + .54 * s, cy + .62 * s),
        (cx + .72 * s, cy - .10 * s),
    ]
    a.draw.polygon([tuple(sc(v) for v in p) for p in pts], fill=rgba(color, alpha), outline=rgba(WHITE, 225))
    a.draw.ellipse(sbox((cx - .22 * s, cy - .52 * s, cx + .06 * s, cy - .22 * s)), fill=rgba(WHITE, 150))


def draw_bubble(a: Art, cx, cy, r, alpha=120):
    a.draw.ellipse(sbox((cx - r, cy - r, cx + r, cy + r)), fill=(255, 255, 255, alpha), outline=(255, 255, 255, min(230, alpha + 70)), width=sc(2))
    a.draw.ellipse(sbox((cx - .38 * r, cy - .38 * r, cx - .08 * r, cy - .08 * r)), fill=(255, 255, 255, 160))


def draw_germ(a: Art, cx, cy, r, color=MINT, mood="smile"):
    edge = mix(color, NAVY, .35)
    for i in range(12):
        angle = 2 * pi * i / 12
        x1 = cx + cos(angle) * r * .72
        y1 = cy + sin(angle) * r * .72
        x2 = cx + cos(angle) * r * 1.18
        y2 = cy + sin(angle) * r * 1.18
        a.draw.line(sbox((x1, y1, x2, y2)), fill=rgba(edge, 210), width=sc(max(3, r * .08)))
        a.draw.ellipse(sbox((x2 - r * .12, y2 - r * .12, x2 + r * .12, y2 + r * .12)), fill=rgba(color, 235), outline=rgba(WHITE, 160), width=sc(2))
    a.circle(cx, cy, r, fill=rgba(color, 235), outline=rgba(WHITE, 190), width=3, shadow=True)
    for ex in (cx - r * .32, cx + r * .32):
        a.draw.ellipse(sbox((ex - r * .10, cy - r * .2, ex + r * .10, cy)), fill=rgba(NAVY, 230))
        a.draw.ellipse(sbox((ex - r * .04, cy - r * .17, ex, cy - r * .12)), fill=rgba(WHITE, 230))
    if mood == "sad":
        a.draw.arc(sbox((cx - r * .34, cy + r * .22, cx + r * .34, cy + r * .78)), 205, 335, fill=rgba(NAVY, 220), width=sc(4))
    elif mood == "wow":
        a.draw.ellipse(sbox((cx - r * .14, cy + r * .18, cx + r * .14, cy + r * .46)), fill=rgba(NAVY, 220))
    else:
        a.draw.arc(sbox((cx - r * .34, cy - r * .05, cx + r * .34, cy + r * .45)), 20, 160, fill=rgba(NAVY, 220), width=sc(4))


def draw_hand(a: Art, cx, cy, scale=1.0, glow=True):
    s = scale
    if glow:
        layer = Image.new("RGBA", a.img.size, (0, 0, 0, 0))
        d = ImageDraw.Draw(layer, "RGBA")
        d.ellipse(sbox((cx - 185 * s, cy - 250 * s, cx + 185 * s, cy + 190 * s)), fill=(255, 214, 155, 80))
        layer = layer.filter(ImageFilter.GaussianBlur(sc(26)))
        a.img.alpha_composite(layer)
        a.draw = ImageDraw.Draw(a.img, "RGBA")
    outline = rgba((178, 108, 84), 210)
    a.draw.rounded_rectangle(sbox((cx - 82 * s, cy - 16 * s, cx + 82 * s, cy + 142 * s)), radius=sc(48 * s), fill=rgba(SKIN, 255), outline=outline, width=sc(4 * s))
    for dx, h in [(-62, 132), (-20, 164), (22, 150), (62, 118)]:
        a.draw.rounded_rectangle(
            sbox((cx + (dx - 19) * s, cy - h * s, cx + (dx + 19) * s, cy + 28 * s)),
            radius=sc(20 * s),
            fill=rgba(SKIN, 255),
            outline=outline,
            width=sc(4 * s),
        )
    a.draw.rounded_rectangle(sbox((cx - 138 * s, cy + 36 * s, cx - 62 * s, cy + 88 * s)), radius=sc(28 * s), fill=rgba(SKIN, 255), outline=outline, width=sc(4 * s))
    for dx in [-62, -20, 22, 62]:
        a.draw.arc(sbox((cx + (dx - 11) * s, cy + 0 * s, cx + (dx + 11) * s, cy + 30 * s)), 0, 180, fill=rgba(SKIN_DARK, 170), width=sc(3 * s))


def draw_two_hands(a: Art, cx, cy, scale=1.0):
    draw_hand(a, cx - 85 * scale, cy, scale * .82)
    draw_hand(a, cx + 95 * scale, cy + 12 * scale, scale * .82)
    for x, y, r in [(cx - 30 * scale, cy - 110 * scale, 38 * scale), (cx + 95 * scale, cy - 125 * scale, 54 * scale), (cx + 180 * scale, cy - 40 * scale, 28 * scale), (cx - 160 * scale, cy - 30 * scale, 32 * scale)]:
        draw_bubble(a, x, y, r, 130)


def draw_soap(a: Art, cx, cy, scale=1.0, color=SOAP):
    s = scale
    a.rounded((cx - 128 * s, cy - 68 * s, cx + 128 * s, cy + 68 * s), radius=45 * s, fill=rgba(color, 230), outline=rgba(WHITE, 210), width=3)
    a.draw.arc(sbox((cx - 70 * s, cy - 24 * s, cx + 70 * s, cy + 42 * s)), 10, 170, fill=rgba(WHITE, 200), width=sc(8 * s))
    for dx, dy, r in [(-90, -78, 22), (96, -85, 30), (20, -100, 16), (-20, 86, 16)]:
        draw_bubble(a, cx + dx * s, cy + dy * s, r * s, 140)


def draw_faucet(a: Art, x, y, scale=1.0):
    s = scale
    metal, edge = (177, 208, 227), (75, 122, 156)
    a.rounded((x, y, x + 285 * s, y + 60 * s), radius=30 * s, fill=rgba(metal, 238), outline=rgba(edge, 230), width=4, shadow=True)
    a.rounded((x + 205 * s, y + 38 * s, x + 265 * s, y + 142 * s), radius=18 * s, fill=rgba(metal, 238), outline=rgba(edge, 230), width=4, shadow=False)
    a.rounded((x + 42 * s, y - 64 * s, x + 106 * s, y + 20 * s), radius=18 * s, fill=rgba(metal, 238), outline=rgba(edge, 230), width=4, shadow=False)
    a.draw.ellipse(sbox((x + 22 * s, y - 95 * s, x + 126 * s, y - 42 * s)), fill=rgba(CORAL, 230), outline=rgba(edge, 230), width=sc(4 * s))
    for i in range(5):
        xx = x + (220 + i * 13) * s
        a.draw.line(sbox((xx, y + 140 * s, xx - 18 * s, y + 345 * s)), fill=rgba(CYAN, 205), width=sc(8 * s))
        draw_drop(a, xx - 24 * s, y + (178 + i * 35) * s, 13 * s, CYAN, 220)


def draw_child(a: Art, cx, cy, scale=1.0, shirt=MINT, accent=GOLD):
    s = scale
    ink = rgba(NAVY, 225)
    # shadow
    a.draw.ellipse(sbox((cx - 115 * s, cy + 205 * s, cx + 115 * s, cy + 248 * s)), fill=(35, 72, 105, 42))
    # legs
    for dx in [-42, 28]:
        a.draw.rounded_rectangle(sbox((cx + dx * s, cy + 96 * s, cx + (dx + 34) * s, cy + 205 * s)), radius=sc(14 * s), fill=rgba((65, 133, 211), 255), outline=ink, width=sc(3 * s))
    a.draw.ellipse(sbox((cx - 76 * s, cy + 195 * s, cx - 5 * s, cy + 228 * s)), fill=rgba(WHITE, 255), outline=ink, width=sc(3 * s))
    a.draw.ellipse(sbox((cx + 8 * s, cy + 195 * s, cx + 82 * s, cy + 228 * s)), fill=rgba(WHITE, 255), outline=ink, width=sc(3 * s))
    # body and arms
    a.draw.rounded_rectangle(sbox((cx - 78 * s, cy - 8 * s, cx + 78 * s, cy + 125 * s)), radius=sc(42 * s), fill=rgba(shirt, 255), outline=ink, width=sc(4 * s))
    a.draw.rounded_rectangle(sbox((cx - 46 * s, cy + 18 * s, cx + 48 * s, cy + 70 * s)), radius=sc(18 * s), fill=rgba(WHITE, 90), outline=rgba(WHITE, 90), width=sc(1))
    a.draw.line(sbox((cx - 68 * s, cy + 42 * s, cx - 132 * s, cy + 103 * s)), fill=rgba(SKIN, 255), width=sc(23 * s))
    a.draw.line(sbox((cx + 68 * s, cy + 42 * s, cx + 132 * s, cy - 8 * s)), fill=rgba(SKIN, 255), width=sc(23 * s))
    a.draw.ellipse(sbox((cx - 152 * s, cy + 86 * s, cx - 116 * s, cy + 124 * s)), fill=rgba(SKIN, 255), outline=ink, width=sc(2 * s))
    a.draw.ellipse(sbox((cx + 116 * s, cy - 28 * s, cx + 154 * s, cy + 10 * s)), fill=rgba(SKIN, 255), outline=ink, width=sc(2 * s))
    # head
    a.draw.ellipse(sbox((cx - 78 * s, cy - 142 * s, cx + 78 * s, cy + 18 * s)), fill=rgba(SKIN, 255), outline=ink, width=sc(4 * s))
    a.draw.pieslice(sbox((cx - 82 * s, cy - 155 * s, cx + 82 * s, cy - 36 * s)), 185, 355, fill=rgba((75, 54, 44), 255), outline=rgba((75, 54, 44), 255))
    a.draw.ellipse(sbox((cx - 85 * s, cy - 92 * s, cx - 58 * s, cy - 48 * s)), fill=rgba((75, 54, 44), 255))
    for ex in [cx - 30 * s, cx + 30 * s]:
        a.draw.ellipse(sbox((ex - 7 * s, cy - 64 * s, ex + 7 * s, cy - 48 * s)), fill=ink)
    a.draw.arc(sbox((cx - 35 * s, cy - 48 * s, cx + 35 * s, cy - 2 * s)), 18, 162, fill=rgba(CORAL, 230), width=sc(4 * s))
    a.draw.ellipse(sbox((cx - 54 * s, cy - 35 * s, cx - 35 * s, cy - 17 * s)), fill=rgba(PINK, 160))
    a.draw.ellipse(sbox((cx + 35 * s, cy - 35 * s, cx + 54 * s, cy - 17 * s)), fill=rgba(PINK, 160))
    draw_sparkle(a, cx + 95 * s, cy - 86 * s, 22 * s, accent)


def draw_shield(a: Art, cx, cy, scale=1.0):
    s = scale
    pts = [
        (cx, cy - 220 * s),
        (cx + 220 * s, cy - 128 * s),
        (cx + 185 * s, cy + 118 * s),
        (cx, cy + 260 * s),
        (cx - 185 * s, cy + 118 * s),
        (cx - 220 * s, cy - 128 * s),
    ]
    a.shadow_round((cx - 240 * s, cy - 220 * s, cx + 240 * s, cy + 285 * s), radius=80 * s, alpha=65)
    a.draw.polygon([tuple(sc(v) for v in p) for p in pts], fill=rgba((255, 237, 130), 235), outline=rgba(WHITE, 230))
    a.draw.line([tuple(sc(v) for v in p) for p in pts + [pts[0]]], fill=rgba(BLUE, 235), width=sc(10 * s))
    a.draw.polygon([tuple(sc(v) for v in p) for p in star_points(cx, cy - 25 * s, 82 * s, 36 * s)], fill=rgba(CORAL, 230), outline=rgba(WHITE, 230))


def draw_clock(a: Art, cx, cy, r):
    a.circle(cx, cy, r, fill=rgba(WHITE, 210), outline=rgba(BLUE, 230), width=8, shadow=True)
    for angle in range(0, 360, 30):
        x1 = cx + cos(angle * pi / 180) * r * .78
        y1 = cy + sin(angle * pi / 180) * r * .78
        x2 = cx + cos(angle * pi / 180) * r * .9
        y2 = cy + sin(angle * pi / 180) * r * .9
        a.draw.line(sbox((x1, y1, x2, y2)), fill=rgba(BLUE, 215), width=sc(4))
    a.draw.line(sbox((cx, cy, cx, cy - r * .55)), fill=rgba(CORAL, 235), width=sc(12))
    a.draw.line(sbox((cx, cy, cx + r * .48, cy + r * .32)), fill=rgba(CORAL, 235), width=sc(12))
    a.circle(cx, cy, 15, fill=rgba(CORAL, 255), shadow=False)


def draw_mini_icon(a: Art, cx, cy, kind, color):
    a.circle(cx, cy, 92, fill=rgba(WHITE, 185), outline=rgba(color, 225), width=5, shadow=True)
    if kind == "meal":
        a.draw.ellipse(sbox((cx - 48, cy - 28, cx + 48, cy + 34)), fill=rgba(WHITE, 255), outline=rgba(NAVY, 220), width=sc(3))
        a.draw.ellipse(sbox((cx - 22, cy - 12, cx + 24, cy + 28)), fill=rgba(GOLD, 255), outline=rgba(NAVY, 180), width=sc(2))
        a.draw.line(sbox((cx + 70, cy - 48, cx + 70, cy + 48)), fill=rgba(NAVY, 210), width=sc(5))
    elif kind == "toilet":
        a.rounded((cx - 40, cy - 52, cx + 36, cy + 4), radius=15, fill=rgba((216, 241, 255), 255), outline=rgba(NAVY, 210), width=3, shadow=False)
        a.rounded((cx - 22, cy + 0, cx + 58, cy + 55), radius=25, fill=rgba(WHITE, 245), outline=rgba(NAVY, 210), width=3, shadow=False)
    elif kind == "sneeze":
        draw_drop(a, cx + 12, cy - 2, 18, CYAN)
        draw_drop(a, cx + 52, cy - 25, 13, CYAN)
        a.draw.arc(sbox((cx - 64, cy - 30, cx + 10, cy + 38)), 260, 75, fill=rgba(CORAL, 230), width=sc(9))
    elif kind == "play":
        a.draw.arc(sbox((cx - 56, cy - 30, cx + 56, cy + 85)), 185, 355, fill=rgba(BLUE, 230), width=sc(8))
        a.draw.line(sbox((cx - 56, cy + 28, cx - 56, cy + 76)), fill=rgba(BLUE, 230), width=sc(8))
        a.draw.line(sbox((cx + 56, cy + 28, cx + 56, cy + 76)), fill=rgba(BLUE, 230), width=sc(8))
        a.circle(cx, cy + 10, 34, fill=rgba(GOLD, 250), outline=rgba(NAVY, 210), width=3, shadow=False)
    elif kind == "pet":
        a.draw.ellipse(sbox((cx - 42, cy - 10, cx + 42, cy + 62)), fill=rgba(WHITE, 255), outline=rgba(NAVY, 210), width=sc(3))
        for dx, dy in [(-45, -25), (-15, -45), (17, -45), (47, -24)]:
            a.draw.ellipse(sbox((cx + dx - 16, cy + dy - 16, cx + dx + 16, cy + dy + 16)), fill=rgba(WHITE, 255), outline=rgba(NAVY, 210), width=sc(3))
    elif kind == "home":
        a.draw.polygon([tuple(sc(v) for v in p) for p in [(cx - 62, cy - 6), (cx, cy - 62), (cx + 62, cy - 6)]], fill=rgba(CORAL, 240), outline=rgba(NAVY, 220))
        a.rounded((cx - 45, cy - 5, cx + 45, cy + 62), radius=10, fill=rgba(WHITE, 235), outline=rgba(NAVY, 220), width=3, shadow=False)


def save_scene(name, scene_func):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    scene_func().save(path)
    return path


def scene_cover():
    a = Art((224, 249, 255), (244, 255, 240), 1)
    a.rounded((230, 160, 1690, 910), radius=74, fill=(255, 255, 255, 86), outline=(255, 255, 255, 170), width=3)
    draw_drop(a, 960, 430, 160, CYAN, 230)
    for r, alpha in [(260, 60), (335, 36), (420, 24)]:
        a.draw.ellipse(sbox((960 - r, 430 - r, 960 + r, 430 + r)), outline=(255, 255, 255, alpha), width=sc(5))
    draw_two_hands(a, 960, 705, 1.15)
    draw_child(a, 445, 680, .9, shirt=(116, 218, 190), accent=GOLD)
    draw_child(a, 1475, 680, .9, shirt=(255, 174, 190), accent=CYAN)
    draw_soap(a, 690, 780, .72, SOAP)
    for x, y, r in [(1240, 245, 42), (1340, 355, 28), (610, 285, 35), (520, 420, 24), (1130, 755, 32)]:
        draw_bubble(a, x, y, r, 125)
    for x, y in [(720, 255), (1210, 585), (360, 280), (1550, 255)]:
        draw_sparkle(a, x, y, 24)
    return a


def scene_busy_hands():
    a = Art((238, 249, 255), (255, 248, 230), 2)
    a.rounded((170, 150, 1750, 850), radius=70, fill=(255, 255, 255, 94), outline=(255, 255, 255, 170), width=3)
    items = [
        (385, 320, "book", MINT),
        (690, 300, "ball", GOLD),
        (995, 300, "handle", CORAL),
        (1300, 320, "apple", PINK),
        (520, 640, "cup", CYAN),
        (870, 650, "toy", LAVENDER),
        (1220, 640, "pet", MINT),
    ]
    for cx, cy, kind, color in items:
        a.circle(cx, cy, 105, fill=rgba(WHITE, 190), outline=rgba(color, 225), width=5)
        if kind == "book":
            a.rounded((cx - 58, cy - 42, cx + 58, cy + 50), radius=12, fill=rgba((108, 179, 255), 235), outline=rgba(NAVY, 220), width=3, shadow=False)
            a.draw.line(sbox((cx, cy - 42, cx, cy + 50)), fill=rgba(WHITE, 200), width=sc(4))
        elif kind == "ball":
            a.circle(cx, cy, 62, fill=rgba(CORAL, 240), outline=rgba(NAVY, 210), width=3, shadow=False)
            a.draw.arc(sbox((cx - 54, cy - 54, cx + 54, cy + 54)), 70, 250, fill=rgba(WHITE, 220), width=sc(5))
        elif kind == "handle":
            a.rounded((cx - 35, cy - 80, cx + 35, cy + 80), radius=16, fill=rgba((204, 161, 110), 240), outline=rgba(NAVY, 210), width=3, shadow=False)
            a.draw.ellipse(sbox((cx - 2, cy - 18, cx + 78, cy + 24)), fill=rgba(GOLD, 245), outline=rgba(NAVY, 210), width=sc(3))
        elif kind == "apple":
            a.draw.ellipse(sbox((cx - 52, cy - 20, cx + 4, cy + 58)), fill=rgba(CORAL, 245), outline=rgba(NAVY, 210), width=sc(3))
            a.draw.ellipse(sbox((cx - 5, cy - 20, cx + 52, cy + 58)), fill=rgba(CORAL, 245), outline=rgba(NAVY, 210), width=sc(3))
            a.draw.line(sbox((cx, cy - 24, cx + 22, cy - 70)), fill=rgba((120, 80, 54), 230), width=sc(6))
        elif kind == "cup":
            a.rounded((cx - 48, cy - 68, cx + 48, cy + 62), radius=20, fill=rgba((157, 224, 255), 230), outline=rgba(NAVY, 210), width=3, shadow=False)
        elif kind == "toy":
            for dx, dy, c in [(-42, 18, PINK), (15, -22, MINT), (58, 22, GOLD)]:
                a.rounded((cx + dx - 37, cy + dy - 37, cx + dx + 37, cy + dy + 37), radius=12, fill=rgba(c, 240), outline=rgba(NAVY, 210), width=3, shadow=False)
        else:
            draw_mini_icon(a, cx, cy, "pet", color)
        draw_germ(a, cx + 84, cy - 78, 22, random.choice([MINT, PINK, LAVENDER, GOLD]), "smile")
    a.draw.line(sbox((305, 505, 1580, 505)), fill=rgba(BLUE, 155), width=sc(10))
    for x in [305, 565, 825, 1085, 1345, 1580]:
        draw_drop(a, x, 505, 18, CYAN, 210)
    return a


def scene_germs():
    a = Art((224, 249, 255), (255, 240, 247), 3)
    a.rounded((155, 150, 760, 835), radius=68, fill=(255, 255, 255, 98), outline=(255, 255, 255, 170), width=3)
    draw_child(a, 445, 675, .92, shirt=GOLD, accent=CYAN)
    draw_hand(a, 620, 520, 1.2)
    a.circle(1200, 505, 310, fill=rgba((237, 253, 255), 184), outline=rgba(BLUE, 230), width=10)
    a.draw.line(sbox((1395, 720, 1660, 925)), fill=rgba(BLUE, 220), width=sc(38))
    for x, y, r, c, m in [
        (1100, 390, 56, MINT, "smile"),
        (1250, 470, 45, CORAL, "wow"),
        (1160, 610, 40, PINK, "smile"),
        (1350, 610, 52, LAVENDER, "smile"),
        (1315, 330, 34, GOLD, "wow"),
    ]:
        draw_germ(a, x, y, r, c, m)
    for x, y, r in [(875, 225, 40), (1600, 210, 34), (1530, 770, 28)]:
        draw_bubble(a, x, y, r, 120)
    return a


def scene_shield():
    a = Art((230, 255, 241), (226, 241, 255), 4)
    a.rounded((190, 135, 1730, 875), radius=78, fill=(255, 255, 255, 84), outline=(255, 255, 255, 170), width=3)
    draw_shield(a, 960, 515, 1.05)
    draw_two_hands(a, 960, 630, .82)
    for gx, gy, r, c in [(350, 330, 56, MINT), (430, 710, 44, PINK), (1545, 360, 58, LAVENDER), (1450, 720, 46, CORAL)]:
        draw_germ(a, gx, gy, r, c, "sad")
        a.draw.line(sbox((gx, gy, 960, 490)), fill=rgba(WHITE, 90), width=sc(5))
    draw_soap(a, 360, 555, .65, SOAP)
    draw_drop(a, 1510, 535, 52, CYAN, 230)
    for x, y in [(620, 230), (1280, 240), (1185, 800), (590, 790)]:
        draw_sparkle(a, x, y, 24)
    return a


def scene_wash_times():
    a = Art((255, 250, 229), (226, 249, 255), 5)
    a.rounded((175, 138, 1745, 884), radius=78, fill=(255, 255, 255, 88), outline=(255, 255, 255, 170), width=3)
    icons = [
        (430, 320, "meal", GOLD),
        (740, 285, "toilet", CYAN),
        (1085, 285, "sneeze", PINK),
        (1400, 325, "play", LAVENDER),
        (590, 665, "pet", MINT),
        (1135, 665, "home", CORAL),
    ]
    for item in icons:
        draw_mini_icon(a, *item)
    draw_child(a, 960, 600, .78, shirt=(116, 218, 190), accent=GOLD)
    return a


def scene_steps():
    a = Art((234, 249, 255), (248, 240, 255), 6)
    points = [(360, 300), (640, 240), (940, 270), (1235, 345), (1370, 680), (985, 790), (565, 690)]
    for i in range(len(points) - 1):
        a.draw.line(sbox((*points[i], *points[i + 1])), fill=rgba(BLUE, 120), width=sc(10))
    for i, (cx, cy) in enumerate(points, 1):
        color = [MINT, GOLD, CYAN, PINK, LAVENDER, CORAL, MINT][i - 1]
        a.circle(cx, cy, 115, fill=rgba(WHITE, 190), outline=rgba(color, 235), width=5)
        draw_two_hands(a, cx, cy + 28, .32)
        a.circle(cx - 72, cy - 72, 30, fill=rgba(color, 250), outline=rgba(WHITE, 230), width=2, shadow=False)
        # The number is only decorative; full Chinese labels are added in PPT.
        for k in range(i):
            draw_bubble(a, cx + 75 - k * 18, cy - 72 + k * 9, 9, 125)
    draw_soap(a, 1480, 775, .64, SOAP)
    return a


def scene_20_seconds():
    a = Art((255, 248, 226), (225, 249, 255), 7)
    draw_clock(a, 960, 445, 250)
    a.rounded((240, 650, 1680, 855), radius=64, fill=(255, 255, 255, 100), outline=(255, 255, 255, 175), width=3)
    draw_hand(a, 445, 735, .65)
    draw_soap(a, 675, 790, .60, SOAP)
    draw_faucet(a, 1225, 670, .62)
    for x, y, r in [(360, 245, 34), (455, 330, 22), (1410, 275, 38), (1510, 370, 26), (1170, 205, 20)]:
        draw_bubble(a, x, y, r, 120)
    for mx, my in [(650, 260), (1300, 560), (510, 545), (1335, 215)]:
        a.draw.ellipse(sbox((mx - 18, my - 18, mx + 18, my + 18)), fill=rgba(LAVENDER, 230), outline=rgba(WHITE, 210), width=sc(2))
        a.draw.line(sbox((mx + 16, my - 8, mx + 16, my - 90)), fill=rgba(LAVENDER, 230), width=sc(7))
    return a


def scene_experiment():
    a = Art((230, 255, 238), (255, 248, 227), 8)
    a.rounded((165, 690, 1755, 900), radius=62, fill=rgba((255, 217, 152), 225), outline=rgba(WHITE, 190), width=3)
    a.circle(750, 485, 255, fill=rgba((232, 249, 255), 190), outline=rgba(BLUE, 230), width=8)
    a.draw.ellipse(sbox((560, 370, 940, 600)), fill=rgba((134, 216, 255), 205), outline=rgba(WHITE, 210), width=sc(4))
    random.seed(82)
    for _ in range(90):
        x = random.randint(585, 915)
        y = random.randint(390, 575)
        if ((x - 750) / 195) ** 2 + ((y - 485) / 110) ** 2 < 1:
            a.draw.ellipse(sbox((x - 4, y - 4, x + 4, y + 4)), fill=rgba(NAVY, 160))
    for angle in range(0, 360, 25):
        x1 = 750 + cos(angle * pi / 180) * 75
        y1 = 485 + sin(angle * pi / 180) * 42
        x2 = 750 + cos(angle * pi / 180) * 170
        y2 = 485 + sin(angle * pi / 180) * 98
        a.draw.line(sbox((x1, y1, x2, y2)), fill=rgba(WHITE, 130), width=sc(4))
    draw_soap(a, 1275, 520, .82, PINK)
    draw_hand(a, 1275, 315, .64)
    draw_child(a, 310, 625, .78, shirt=MINT, accent=GOLD)
    return a


def scene_cough():
    a = Art((255, 240, 240), (228, 249, 255), 9)
    draw_child(a, 575, 650, 1.02, shirt=GOLD, accent=CYAN)
    a.rounded((720, 455, 900, 560), radius=45, fill=rgba(SKIN, 245), outline=rgba(WHITE, 185), width=3)
    a.rounded((870, 430, 1040, 575), radius=28, fill=rgba(WHITE, 230), outline=rgba(BLUE, 220), width=4)
    a.rounded((1070, 250, 1160, 720), radius=40, fill=(255, 255, 255, 150), outline=rgba(CORAL, 230), width=8)
    for yy in [335, 485, 635]:
        a.draw.line(sbox((1090, yy, 1140, yy)), fill=rgba(CORAL, 230), width=sc(12))
    for x, y, r, c in [(1260, 310, 44, MINT), (1420, 460, 40, PINK), (1290, 650, 36, LAVENDER)]:
        draw_germ(a, x, y, r, c, "wow")
    a.rounded((1370, 720, 1555, 910), radius=28, fill=rgba((178, 225, 197), 240), outline=rgba(NAVY, 220), width=4)
    a.draw.rectangle(sbox((1348, 682, 1578, 730)), fill=rgba((108, 168, 136), 245), outline=rgba(NAVY, 220), width=sc(3))
    return a


def scene_clean_habits():
    a = Art((244, 241, 255), (224, 251, 242), 10)
    a.rounded((180, 150, 1740, 865), radius=78, fill=(255, 255, 255, 92), outline=(255, 255, 255, 170), width=3)
    draw_hand(a, 455, 540, 1.05)
    a.rounded((700, 310, 1035, 395), radius=25, fill=rgba((184, 209, 224), 240), outline=rgba(NAVY, 220), width=4)
    a.rounded((780, 245, 1060, 310), radius=18, fill=rgba((225, 239, 248), 245), outline=rgba(NAVY, 220), width=4)
    a.draw.line(sbox((1000, 270, 1090, 190)), fill=rgba(NAVY, 220), width=sc(8))
    a.rounded((790, 570, 1145, 700), radius=52, fill=rgba((255, 210, 171), 245), outline=rgba(NAVY, 220), width=4)
    a.rounded((935, 598, 1005, 672), radius=15, fill=rgba((255, 239, 211), 245), outline=rgba(SKIN_DARK, 220), width=2)
    a.rounded((1375, 245, 1530, 560), radius=34, fill=rgba((137, 219, 255), 235), outline=rgba(NAVY, 220), width=4)
    a.rounded((1350, 650, 1565, 850), radius=44, fill=rgba(WHITE, 235), outline=rgba(PINK, 230), width=5)
    for yy in [690, 735, 780, 825]:
        a.draw.line(sbox((1395, yy, 1525, yy)), fill=rgba(PINK, 185), width=sc(7))
    for x, y in [(250, 260), (1215, 205), (1575, 350), (1220, 800)]:
        draw_sparkle(a, x, y, 26)
    return a


def scene_detective():
    a = Art((232, 249, 255), (255, 248, 230), 11)
    a.rounded((165, 145, 1080, 850), radius=62, fill=rgba(WHITE, 170), outline=rgba(BLUE, 210), width=5)
    rooms = [
        (230, 220, MINT),
        (505, 220, GOLD),
        (780, 220, PINK),
        (230, 510, CORAL),
        (505, 510, CYAN),
        (780, 510, LAVENDER),
    ]
    for x, y, color in rooms:
        a.rounded((x, y, x + 205, y + 180), radius=26, fill=rgba(color, 225), outline=rgba(WHITE, 200), width=3, shadow=False)
        draw_drop(a, x + 103, y + 90, 28, WHITE, 225)
    a.draw.line(sbox((310, 755, 945, 755)), fill=rgba(GOLD, 190), width=sc(16))
    for x in [310, 485, 660, 835, 945]:
        a.circle(x, 755, 18, fill=rgba(GOLD, 255), outline=rgba(WHITE, 220), width=2, shadow=False)
    a.rounded((1280, 210, 1585, 645), radius=35, fill=rgba(CREAM, 245), outline=rgba(NAVY, 220), width=5)
    a.rounded((1370, 175, 1495, 245), radius=22, fill=rgba(CORAL, 240), outline=rgba(NAVY, 220), width=4)
    for i in range(5):
        yy = 305 + i * 60
        a.rounded((1320, yy - 16, 1350, yy + 16), radius=8, fill=rgba(MINT, 240), outline=rgba(NAVY, 220), width=2, shadow=False)
        a.draw.line(sbox((1380, yy, 1540, yy)), fill=rgba(NAVY, 115), width=sc(4))
    draw_child(a, 1225, 775, .68, shirt=PINK, accent=GOLD)
    a.circle(1180, 245, 74, fill=rgba((230, 255, 255), 185), outline=rgba(BLUE, 220), width=5)
    a.draw.line(sbox((1230, 300, 1320, 385)), fill=rgba(BLUE, 220), width=sc(18))
    return a


def scene_pledge():
    a = Art((219, 246, 255), (255, 249, 225), 12)
    a.rounded((205, 160, 1715, 895), radius=82, fill=(255, 255, 255, 86), outline=(255, 255, 255, 175), width=3)
    a.rounded((250, 745, 1670, 925), radius=55, fill=rgba(WHITE, 210), outline=rgba(BLUE, 180), width=4)
    draw_child(a, 525, 650, .92, shirt=PINK, accent=GOLD)
    draw_child(a, 960, 625, 1.05, shirt=GOLD, accent=CYAN)
    draw_child(a, 1395, 650, .92, shirt=MINT, accent=GOLD)
    for cx, cy, color in [(525, 290, PINK), (960, 250, GOLD), (1395, 290, MINT)]:
        a.draw.polygon([tuple(sc(v) for v in p) for p in star_points(cx, cy, 78, 35)], fill=rgba(color, 235), outline=rgba(WHITE, 230))
        draw_drop(a, cx, cy, 26, CYAN, 225)
    draw_soap(a, 300, 670, .62, LAVENDER)
    draw_faucet(a, 1515, 625, .50)
    for x, y, r in [(255, 265, 42), (1545, 260, 42), (760, 850, 32), (1160, 820, 28)]:
        draw_bubble(a, x, y, r, 125)
    return a


SCENES = [
    ("slide01_premium_cover.png", scene_cover),
    ("slide02_premium_busy_hands.png", scene_busy_hands),
    ("slide03_premium_germs.png", scene_germs),
    ("slide04_premium_shield.png", scene_shield),
    ("slide05_premium_wash_times.png", scene_wash_times),
    ("slide06_premium_steps.png", scene_steps),
    ("slide07_premium_20_seconds.png", scene_20_seconds),
    ("slide08_premium_experiment.png", scene_experiment),
    ("slide09_premium_cough.png", scene_cough),
    ("slide10_premium_clean_habits.png", scene_clean_habits),
    ("slide11_premium_detective.png", scene_detective),
    ("slide12_premium_pledge.png", scene_pledge),
]


def add_bg(slide, image_path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(SLIDE_W), height=Inches(SLIDE_H))


def style_run(run, size=22, bold=False, color=NAVY):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def add_box(slide, x, y, w, h, fill=WHITE, outline=WHITE, transparency=10, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(*fill)
    box.fill.transparency = transparency
    box.line.color.rgb = RGBColor(*outline)
    box.line.transparency = 15
    box.line.width = Pt(1.5)
    return box


def add_text(slide, x, y, w, h, text, size=22, bold=False, color=NAVY, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for run in p.runs:
        style_run(run, size=size, bold=bold, color=color)
    return tx


def add_title(slide, title, subtitle=None):
    add_box(slide, .45, .28, 12.43, .92, fill=WHITE, outline=(210, 239, 255), transparency=7)
    add_text(slide, .78, .38, 11.78, .42, title, size=26, bold=True, color=(25, 93, 154), align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, 1.15, .82, 11.05, .25, subtitle, size=10.5, color=(74, 104, 133), align=PP_ALIGN.CENTER)


def add_footer(slide, page):
    add_box(slide, 11.95, 7.03, .88, .30, fill=WHITE, outline=(198, 232, 255), transparency=12)
    add_text(slide, 12.07, 7.07, .62, .17, f"{page:02d}", size=9, bold=True, color=(25, 93, 154), align=PP_ALIGN.CENTER)


def add_badge(slide, x, y, w, h, text, fill=(255, 246, 191), color=NAVY, size=14):
    add_box(slide, x, y, w, h, fill=fill, outline=WHITE, transparency=0)
    add_text(slide, x + .08, y + .05, w - .16, h - .10, text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_bullets(slide, x, y, w, h, bullets, size=16, fill=WHITE):
    add_box(slide, x, y, w, h, fill=fill, outline=WHITE, transparency=8)
    tx = slide.shapes.add_textbox(Inches(x + .22), Inches(y + .16), Inches(w - .42), Inches(h - .25))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, text in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = "• " + text
        p.space_after = Pt(5)
        p.line_spacing = 1.12
        for run in p.runs:
            style_run(run, size=size, color=(45, 74, 103))


def build_ppt(image_paths):
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]
    slides = []
    for image_path in image_paths:
        slide = prs.slides.add_slide(blank)
        add_bg(slide, image_path)
        slides.append(slide)

    add_box(slides[0], 2.2, .72, 8.92, 1.55, fill=WHITE, outline=(214, 240, 255), transparency=4)
    add_text(slides[0], 2.55, .88, 8.25, .58, "小学生手卫生", size=39, bold=True, color=(24, 92, 154), align=PP_ALIGN.CENTER)
    add_text(slides[0], 2.85, 1.56, 7.65, .30, "二年级健康科普课 | 和泡泡一起守护干净小手", size=14, bold=True, color=(224, 100, 87), align=PP_ALIGN.CENTER)
    add_badge(slides[0], 4.86, 6.55, 3.62, .48, "伸出小手，开始健康冒险", fill=(255, 239, 166), size=15)
    add_footer(slides[0], 1)

    add_title(slides[1], "小手每天很忙")
    add_bullets(slides[1], .62, 1.45, 4.15, 2.34, [
        "写字、玩耍、摸门把手，小手到处帮忙。",
        "看不见的小细菌，也可能悄悄搭便车。",
        "学会洗手，小手就能清爽又可靠。"
    ], size=16.5)
    add_badge(slides[1], 4.55, 6.66, 4.25, .46, "互动：今天你的小手碰过什么？", fill=(255, 239, 166), size=14.5)
    add_footer(slides[1], 2)

    add_title(slides[2], "细菌病毒在哪里？")
    add_bullets(slides[2], .62, 1.42, 4.08, 2.62, [
        "它们很小很小，肉眼看不见。",
        "喜欢躲在手指缝、指甲边和手心里。",
        "有些会让我们肚子疼、咳嗽或发烧。"
    ], size=16.5)
    add_badge(slides[2], .82, 4.30, 3.55, .48, "脏手要洗干净，健康更安心", fill=(218, 250, 239), size=13.5)
    add_footer(slides[2], 3)

    add_title(slides[3], "洗手像给身体加护盾")
    add_bullets(slides[3], .72, 1.45, 4.02, 2.45, [
        "洗掉灰尘、油油和脏东西。",
        "减少细菌从手跑到口、眼、鼻。",
        "保护自己，也保护同学和家人。"
    ], size=16.5)
    add_badge(slides[3], 9.05, 6.42, 3.45, .50, "干净小手 = 健康加分", fill=(255, 239, 166), size=16)
    add_footer(slides[3], 4)

    add_title(slides[4], "这些时候一定要洗手")
    for text, x, y, color in [
        ("吃东西前", 2.15, 2.65, (255, 239, 166)),
        ("上厕所后", 4.28, 2.42, (214, 244, 255)),
        ("咳嗽喷嚏后", 6.55, 2.42, (255, 222, 229)),
        ("户外玩耍后", 8.62, 2.76, (231, 224, 255)),
        ("摸宠物后", 3.18, 5.38, (218, 250, 239)),
        ("回家或进教室后", 6.75, 5.40, (255, 230, 204)),
    ]:
        add_badge(slides[4], x, y, 1.72, .42, text, fill=color, size=12.5)
    add_footer(slides[4], 5)

    add_title(slides[5], "七步洗手法：小手面面俱到")
    add_bullets(slides[5], 9.05, 1.32, 3.62, 4.75, [
        "1 掌心相对搓一搓",
        "2 手背交叉搓一搓",
        "3 手指交叉搓一搓",
        "4 弯弯指背搓一搓",
        "5 大拇指转一转",
        "6 指尖掌心画圈圈",
        "7 手腕也要洗干净",
    ], size=13.5)
    add_badge(slides[5], 4.34, 6.92, 4.72, .38, "每一步都搓到，细菌没处躲", fill=(255, 239, 166), size=13)
    add_footer(slides[5], 6)

    add_title(slides[6], "洗手小口诀")
    for i, (text, color) in enumerate([
        ("湿一湿", (214, 244, 255)),
        ("抹香皂", (231, 224, 255)),
        ("搓20秒", (255, 239, 166)),
        ("冲干净", (218, 250, 239)),
        ("擦干手", (255, 222, 229)),
    ]):
        add_badge(slides[6], .72, 1.38 + i * 1.05, 2.08, .62, text, fill=color, size=18)
    add_bullets(slides[6], 9.12, 1.50, 3.40, 2.25, [
        "可以唱一遍生日歌。",
        "手心、手背、指缝都要有泡泡。",
        "擦干后，小手不再湿哒哒。"
    ], size=15)
    add_badge(slides[6], 9.00, 5.88, 3.82, .50, "20秒，让泡泡认真工作", fill=(255, 239, 166), size=15)
    add_footer(slides[6], 7)

    add_title(slides[7], "泡泡科学小实验")
    add_bullets(slides[7], 9.08, 1.30, 3.58, 3.44, [
        "盘里倒清水，撒一点胡椒粉。",
        "手指先碰清水，胡椒还在附近。",
        "手指沾肥皂再碰，胡椒会散开。",
        "在老师或家长陪同下做实验。"
    ], size=14.5)
    add_badge(slides[7], 4.55, 6.60, 4.34, .48, "肥皂能帮水带走油油和脏东西", fill=(218, 250, 239), size=14)
    add_footer(slides[7], 8)

    add_title(slides[8], "打喷嚏也要讲卫生")
    add_bullets(slides[8], .68, 1.45, 3.95, 3.30, [
        "用纸巾或手肘挡住口鼻。",
        "用过的纸巾丢进垃圾桶。",
        "咳嗽、打喷嚏后要洗手。",
        "别用脏手揉眼睛、抠鼻子。"
    ], size=15.5)
    add_badge(slides[8], 8.42, 6.47, 3.85, .48, "礼貌一挡，细菌少飞翔", fill=(255, 239, 166), size=15)
    add_footer(slides[8], 9)

    add_title(slides[9], "让小手保持干净")
    add_bullets(slides[9], .72, 1.36, 3.80, 3.10, [
        "勤剪指甲，别把脏东西藏起来。",
        "不咬手指，不用手乱摸脸。",
        "小伤口贴创可贴，保持干净。",
        "毛巾、水杯等个人物品不混用。"
    ], size=15.2)
    add_badge(slides[9], 7.82, 6.55, 4.15, .48, "干净小习惯，每天都能做", fill=(218, 250, 239), size=15)
    add_footer(slides[9], 10)

    add_title(slides[10], "校园洗手小侦探")
    add_bullets(slides[10], 8.50, 1.25, 3.95, 1.14, [
        "找一找：学校里哪里可以洗手？",
        "记一记：今天抓住几个洗手机会？"
    ], size=13.5)
    add_badge(slides[10], 8.95, 6.10, 3.35, .48, "完成一次，给自己一颗星", fill=(255, 239, 166), size=14)
    add_badge(slides[10], 2.08, 6.93, 4.78, .34, "课堂互动：做一名“洗手提醒员”", fill=(214, 244, 255), size=12)
    add_footer(slides[10], 11)

    add_title(slides[11], "我是手卫生小英雄")
    add_box(slides[11], 2.18, 1.24, 8.98, 1.45, fill=WHITE, outline=(214, 240, 255), transparency=4)
    add_text(slides[11], 2.60, 1.38, 8.18, .34, "我承诺：", size=22, bold=True, color=(224, 100, 87), align=PP_ALIGN.CENTER)
    add_text(slides[11], 2.64, 1.88, 8.10, .42, "吃前便后洗手，搓够20秒，提醒伙伴一起做！", size=18, bold=True, color=(24, 92, 154), align=PP_ALIGN.CENTER)
    add_badge(slides[11], 4.16, 6.54, 5.04, .50, "小手干净，健康同行！", fill=(255, 239, 166), size=19)
    add_footer(slides[11], 12)

    prs.save(OUT_FILE)


def main():
    old_assets = list(ASSET_DIR.glob("slide*.png")) if ASSET_DIR.exists() else []
    for old_asset in old_assets:
        old_asset.unlink()
    image_paths = [save_scene(name, func) for name, func in SCENES]
    build_ppt(image_paths)
    print(f"Generated: {OUT_FILE}")
    print(f"Slides: {len(image_paths)}")
    print(f"Assets: {ASSET_DIR}")


if __name__ == "__main__":
    main()
